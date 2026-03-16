from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


TITLE_COLOR = RGBColor(18, 39, 63)
ACCENT_COLOR = RGBColor(214, 98, 46)
TEXT_COLOR = RGBColor(47, 55, 65)
MUTED_COLOR = RGBColor(96, 106, 118)
LIGHT_BG = RGBColor(247, 244, 238)
CODE_BG = RGBColor(243, 246, 249)
WHITE = RGBColor(255, 255, 255)


def add_background(slide, prs):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        prs.slide_width,
        prs.slide_height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.fill.background()
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(11.5), Inches(0.8))
    frame = title_box.text_frame
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(1.1), Inches(11.0), Inches(0.45)
        )
        sub_frame = subtitle_box.text_frame
        p = sub_frame.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        run.font.name = "Aptos"
        run.font.size = Pt(12)
        run.font.color.rgb = MUTED_COLOR


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.75), Inches(6.9), Inches(11.4), Inches(0.25))
    frame = box.text_frame
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(9)
    run.font.color.rgb = MUTED_COLOR


def add_section_card(slide, heading, body, left, top, width, height):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(226, 221, 214)

    heading_box = slide.shapes.add_textbox(
        Inches(left + 0.2), Inches(top + 0.12), Inches(width - 0.35), Inches(0.3)
    )
    heading_frame = heading_box.text_frame
    p = heading_frame.paragraphs[0]
    run = p.add_run()
    run.text = heading
    run.font.name = "Aptos"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = ACCENT_COLOR

    body_box = slide.shapes.add_textbox(
        Inches(left + 0.2), Inches(top + 0.48), Inches(width - 0.35), Inches(height - 0.58)
    )
    body_frame = body_box.text_frame
    body_frame.word_wrap = True
    for index, line in enumerate(body.split("\n")):
        p = body_frame.paragraphs[0] if index == 0 else body_frame.add_paragraph()
        p.text = line
        p.font.name = "Aptos"
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(4)


def add_bullets(slide, items, left, top, width, height, font_size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = item
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(8)


def add_code_box(slide, heading, code, left, top, width, height, font_size=11.5):
    label = slide.shapes.add_textbox(Inches(left), Inches(top - 0.28), Inches(width), Inches(0.25))
    label_frame = label.text_frame
    p = label_frame.paragraphs[0]
    run = p.add_run()
    run.text = heading
    run.font.name = "Aptos"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = ACCENT_COLOR

    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.color.rgb = RGBColor(220, 225, 231)

    box = slide.shapes.add_textbox(
        Inches(left + 0.14), Inches(top + 0.1), Inches(width - 0.28), Inches(height - 0.2)
    )
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    for index, line in enumerate(code.strip().split("\n")):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = line.rstrip()
        p.font.name = "Consolas"
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(0)


def build_presentation(output_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)
    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.75), Inches(0.8), Inches(0.18), Inches(4.8)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_COLOR
    accent.line.fill.background()
    add_title(slide, "Decouverte Angular 21", "Support de presentation pour une equipe habituee a Angular 16")
    add_section_card(
        slide,
        "Objectif",
        "Comprendre les changements utiles au quotidien:\nStandalone, nouveau controle de flux,\nSignals, interop RxJS, Signal Store et Signal Forms.",
        1.2,
        2.8,
        4.8,
        1.95,
    )
    add_section_card(
        slide,
        "Format",
        "Une presentation par chapitres.\nLes exercices servent ensuite a pratiquer chaque sujet.",
        6.2,
        2.8,
        4.8,
        1.95,
    )
    add_footer(slide, "ATOS - Angular 21 workshop")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)
    add_title(slide, "Plan de presentation")
    add_bullets(
        slide,
        [
            "Le fil rouge des evolutions Angular recentes",
            "Standalone",
            "Nouveau controle de flux et @let",
            "Signals",
            "input() et output()",
            "Interop RxJS / Signals",
            "NgRx Signal Store",
            "Signal Forms",
            "Conclusion",
        ],
        0.95,
        1.7,
        10.8,
        4.6,
        20,
    )
    add_footer(slide, "Chapitrage general")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)
    add_title(slide, "Le fil rouge des evolutions Angular recentes")
    add_section_card(
        slide,
        "Avant",
        "Modules de structuration\nTemplates parfois verbeux\nDerivees recalculees a la main\nSeparation moins nette entre etat et effets",
        0.85,
        1.65,
        4.95,
        3.75,
    )
    add_section_card(
        slide,
        "Direction actuelle",
        "Composants plus autonomes\nTemplates plus declaratifs\nState local plus explicite\nInterop officielle entre les modeles reactifs",
        6.0,
        1.65,
        4.95,
        3.75,
    )
    add_bullets(
        slide,
        [
            "Angular ne dit pas que l ancien code est mauvais.",
            "Il pousse surtout un code plus local, plus lisible et plus simple a raisonner.",
        ],
        1.0,
        5.65,
        10.1,
        0.9,
        17,
    )
    add_footer(slide, "Fil directeur")

    standalone_before = """@NgModule({
  declarations: [ProfileFormComponent],
  imports: [
    CommonModule,
    FormsModule,
    ButtonModule,
    InputTextModule,
    RouterModule.forChild(routes)
  ]
})
export class TeamFeatureModule {}"""

    standalone_after = """@Component({
  standalone: true,
  imports: [FormsModule, Button, InputText],
  template: `
    <input pInputText [(ngModel)]="fullName" />
    <button pButton type="button" label="Enregistrer"></button>
  `
})
export class ProfileFormComponent {
  fullName = '';
}"""

    standalone_routes = """{
  path: 'profile',
  loadComponent: () =>
    import('./profile-form.component').then((m) => m.ProfileFormComponent)
}

{
  path: 'admin',
  loadChildren: () =>
    import('./admin/admin.routes').then((m) => m.adminRoutes)
}"""

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)
    add_title(slide, "Chapitre 1 - Standalone", "Simplifier l architecture et localiser les dependances")
    add_section_card(
        slide,
        "Ce que ca change",
        "On raisonne moins en NgModule et davantage en composants autonomes qui declarent eux-memes leurs dependances.",
        0.85,
        1.55,
        3.35,
        2.05,
    )
    add_section_card(
        slide,
        "Impact",
        "On lit plus vite une feature.\nMoins de boilerplate.\nLe lazy loading devient plus direct.",
        4.35,
        1.55,
        3.15,
        2.05,
    )
    add_section_card(
        slide,
        "Utilite",
        "Tres parlant avec PrimeNG:\nle composant importe lui-meme Button, InputText, Select, etc.",
        7.7,
        1.55,
        3.8,
        2.05,
    )
    add_code_box(slide, "Avant", standalone_before, 0.85, 3.95, 3.6, 2.45)
    add_code_box(slide, "Maintenant avec PrimeNG", standalone_after, 4.65, 3.95, 3.75, 2.45)
    add_code_box(slide, "loadComponent / loadChildren", standalone_routes, 8.6, 3.95, 2.85, 2.45, 10.5)
    add_footer(slide, "Message cle: les dependances sont visibles la ou elles servent")

    flow_before = """<section *ngIf="filteredMovies.length > 0; else emptyState">
  <article *ngFor="let movie of filteredMovies; trackBy: trackByMovieId">
    {{ movie.title }}
  </article>
</section>

<ng-template #emptyState>
  <p>Aucun film trouve.</p>
</ng-template>"""

    flow_after = """@let visibleMovies = filteredMovies();

@if (visibleMovies.length > 0) {
  @for (movie of visibleMovies; track movie.id) {
    <article>{{ movie.title }}</article>
  }
} @else {
  <p>Aucun film trouve.</p>
}"""

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)
    add_title(slide, "Chapitre 2 - Nouveau controle de flux", "Des templates plus lisibles avec @if, @for, @switch et @let")
    add_section_card(
        slide,
        "Ce que ca change",
        "Angular propose une syntaxe de template plus declarative pour les conditions, les listes et les valeurs intermediaires.",
        0.85,
        1.55,
        3.4,
        2.0,
    )
    add_section_card(
        slide,
        "Impact",
        "Le template devient plus narratif.\nLe track ressort mieux.\n@let reduit les repetitions.",
        4.45,
        1.55,
        3.0,
        2.0,
    )
    add_section_card(
        slide,
        "Utilite",
        "Tres utile des qu un template a plusieurs branches ou repete les memes expressions.",
        7.7,
        1.55,
        3.8,
        2.0,
    )
    add_code_box(slide, "Avant", flow_before, 0.85, 3.95, 5.15, 2.45)
    add_code_box(slide, "Maintenant", flow_after, 6.15, 3.95, 5.3, 2.45)
    add_footer(slide, "Message cle: la nouvelle syntaxe raconte mieux l intention du template")

    signals_before = """searchTerm = '';
selectedCategory = 'all';
filteredProducts: Product[] = [];
resultsCount = 0;
pageTitle = 'Catalogue';

refreshViewModel(): void {
  this.filteredProducts = this.products.filter(...);
  this.resultsCount = this.filteredProducts.length;
  this.pageTitle = `Catalogue (${this.resultsCount})`;
}"""

    signals_after = """readonly searchTerm = signal('');
readonly selectedCategory = signal('all');
readonly filteredProducts = computed(() => ...);
readonly resultsCount = computed(() => this.filteredProducts().length);
readonly pageTitle = computed(() => `Catalogue (${this.resultsCount()})`);

effect(() => {
  document.title = this.pageTitle();
});"""

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)
    add_title(slide, "Chapitre 3 - Signals", "Repenser le state local avec signal(), computed() et effect()")
    add_section_card(
        slide,
        "Ce que ca change",
        "On distingue beaucoup mieux la source de verite, la derivee et l effet de bord.",
        0.85,
        1.55,
        3.4,
        1.95,
    )
    add_section_card(
        slide,
        "Impact",
        "Moins de derivees recalculees a la main.\nMoins de risque de desynchronisation.\nPlus de lisibilite.",
        4.45,
        1.55,
        3.15,
        1.95,
    )
    add_section_card(
        slide,
        "Utilite",
        "Filtres, compteurs, tri, derivees d affichage et petits view models locaux.",
        7.8,
        1.55,
        3.7,
        1.95,
    )
    add_code_box(slide, "Avant", signals_before, 0.85, 3.85, 5.15, 2.45)
    add_code_box(slide, "Maintenant", signals_after, 6.15, 3.85, 5.3, 2.45)
    add_footer(slide, "A retenir: signal = source, computed = derivee, effect = effet de bord")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)
    add_title(slide, "Chapitre 3 - Definitions simples", "Des exemples courts pour fixer les roles de chaque primitive")
    add_code_box(
        slide,
        "signal()",
        """readonly count = signal(0);
readonly searchTerm = signal('');

this.count.set(1);
this.count.update((value) => value + 1);""",
        0.8,
        1.9,
        3.75,
        1.85,
    )
    add_code_box(
        slide,
        "computed()",
        """readonly price = signal(100);
readonly quantity = signal(2);
readonly total = computed(() =>
  this.price() * this.quantity()
);""",
        4.8,
        1.9,
        3.75,
        1.85,
    )
    add_code_box(
        slide,
        "effect()",
        """readonly pageTitle = computed(() =>
  `Panier (${this.items().length})`
);

effect(() => {
  document.title = this.pageTitle();
});""",
        8.8,
        1.9,
        3.75,
        1.85,
    )
    add_bullets(
        slide,
        [
            "signal() stocke une source de verite locale.",
            "computed() derive une nouvelle valeur a partir d autres signals.",
            "effect() synchronise un effet externe.",
            "A eviter: utiliser effect() pour recalculer un state metier qui pourrait etre un computed().",
        ],
        0.95,
        4.35,
        11.0,
        1.8,
        16,
    )
    add_footer(slide, "Regle simple: source, derivee, effet")

    io_before = """@Input() resultsCount = 0;
@Input() pageTitle = '';
@Output() availableOnlyChange = new EventEmitter<boolean>();"""

    io_after = """readonly resultsCount = input.required<number>();
readonly pageTitle = input.required<string>();
readonly availableOnly = input(false);
readonly availableOnlyChange = output<boolean>();

readonly summaryLabel = computed(() =>
  `${this.pageTitle()} - ${this.resultsCount()} resultat(s)`
);"""

    io_effect = """constructor() {
  effect(() => {
    console.log(
      `Le parent a envoye ${this.resultsCount()} resultat(s)`
    );
  });
}"""

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)
    add_title(slide, "Chapitre 4 - input() et output()", "Les input() s integrent naturellement au modele Signals")
    add_section_card(
        slide,
        "Ce que ca change",
        "Angular propose des APIs plus directes et plus coherentes avec les autres primitives modernes.",
        0.85,
        1.55,
        3.45,
        1.95,
    )
    add_section_card(
        slide,
        "Impact",
        "Declaration plus concise.\nLecture plus simple.\nPossibilite de deriver directement dans l enfant.",
        4.45,
        1.55,
        3.25,
        1.95,
    )
    add_section_card(
        slide,
        "Utilite",
        "Rafraichir un composant legacy sans changer la responsabilite du parent.",
        7.9,
        1.55,
        3.55,
        1.95,
    )
    add_code_box(slide, "Avant", io_before, 0.85, 3.8, 3.45, 1.55)
    add_code_box(slide, "Maintenant", io_after, 4.5, 3.8, 3.6, 1.55)
    add_code_box(slide, "effect() dans l enfant", io_effect, 8.3, 3.8, 3.15, 1.55)
    add_code_box(
        slide,
        "Utilisation cote parent",
        """<app-product-summary
  [resultsCount]="resultsCount()"
  [pageTitle]="pageTitle()"
  [availableOnly]="availableOnly()"
  (availableOnlyChange)="setAvailableOnly($event)"
/>""",
        0.85,
        5.75,
        10.6,
        1.0,
    )
    add_footer(slide, "Message cle: un input() est lisible comme un signal dans l enfant")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)
    add_title(slide, "Chapitre 5 - Promise, Observable, Signal", "Trois modeles differents, trois usages differents")
    add_section_card(
        slide,
        "Promise",
        "Une valeur unique plus tard.\nBonne option pour une operation ponctuelle.",
        0.8,
        1.7,
        3.8,
        1.55,
    )
    add_section_card(
        slide,
        "Observable",
        "Un flux de valeurs dans le temps.\nIdeal pour les streams, events et pipelines RxJS.",
        4.8,
        1.7,
        3.8,
        1.55,
    )
    add_section_card(
        slide,
        "Signal",
        "Une valeur courante reactive.\nTres pratique pour le state local et les derivees UI.",
        8.8,
        1.7,
        3.75,
        1.55,
    )
    add_code_box(
        slide,
        "Promise",
        """async function loadUser(): Promise<User> {
  const response = await fetch('/api/user');
  return response.json();
}""",
        0.8,
        3.65,
        3.8,
        1.55,
    )
    add_code_box(
        slide,
        "Observable",
        """readonly query$ = new BehaviorSubject('');
readonly filteredProducts$ = combineLatest([
  this.products$,
  this.query$
]).pipe(map(...));""",
        4.8,
        3.65,
        3.8,
        1.55,
    )
    add_code_box(
        slide,
        "Signal",
        """readonly query = signal('');
readonly filteredProducts = computed(() =>
  this.products().filter(...)
);""",
        8.8,
        3.65,
        3.75,
        1.55,
    )
    add_bullets(
        slide,
        [
            "Le bon raisonnement n est pas de choisir un vainqueur.",
            "Le bon raisonnement est de choisir l outil adapte au probleme traite.",
        ],
        1.0,
        5.7,
        10.3,
        0.9,
        17,
    )
    add_footer(slide, "Comparer sans opposer")

    interop_before = """readonly filteredProducts$ = combineLatest([
  this.productStream$,
  this.query$
]).pipe(
  map(([products, query]) => filterProducts(products, query))
);"""

    interop_after = """readonly products = toSignal(this.productService.products$, {
  initialValue: []
});
readonly query = signal('');
readonly filteredProducts = computed(() => {
  return this.products().filter((product) =>
    product.name.toLowerCase().includes(this.query().toLowerCase())
  );
});"""

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)
    add_title(slide, "Chapitre 5 - Interop RxJS / Signals", "Garder RxJS dans les services et simplifier la lecture dans le composant")
    add_code_box(slide, "Avant", interop_before, 0.8, 1.9, 5.3, 3.9)
    add_code_box(slide, "Maintenant", interop_after, 6.25, 1.9, 5.3, 3.9)
    add_bullets(
        slide,
        [
            "Le service peut rester en RxJS.",
            "Le composant lit une valeur courante avec toSignal().",
            "On simplifie la couche UI sans casser le modele existant.",
        ],
        1.0,
        6.05,
        10.2,
        0.6,
        16,
    )
    add_footer(slide, "Message cle: toSignal() simplifie la couche UI sans remplacer RxJS")

    ngrx_before = """readonly products = signal<Product[]>([]);
readonly favoriteIds = signal<number[]>([]);
readonly loading = signal(false);
readonly errorMessage = signal<string | null>(null);

async loadProducts(): Promise<void> {
  this.loading.set(true);
  try {
    this.products.set(await fetchProducts());
  } finally {
    this.loading.set(false);
  }
}"""

    ngrx_after = """export const ProductsStore = signalStore(
  { providedIn: 'root' },
  withEntities<Product>(),
  withState({ favoriteIds: [], filter: 'all' }),
  withComputed((store) => ({
    favoriteCount: computed(() => store.favoriteIds().length)
  })),
  withMethods((store) => ({ toggleFavorite(productId: number) { ... } })),
  withHooks({ onInit(store) { void store.loadProducts(); } })
);"""

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)
    add_title(slide, "Chapitre 6 - NgRx Signal Store", "Structurer un state partage et eviter les rechargements inutiles")
    add_section_card(
        slide,
        "Ce que ca change",
        "On sort du simple state local pour centraliser l etat partage entre plusieurs vues.",
        0.85,
        1.55,
        3.4,
        1.95,
    )
    add_section_card(
        slide,
        "Impact",
        "Une seule source de verite.\nMoins d appels redondants.\nUne logique mieux separee.",
        4.45,
        1.55,
        3.1,
        1.95,
    )
    add_section_card(
        slide,
        "Utilite",
        "Catalogue partage, favoris, filtres, statut de chargement, donnees conservees tant que la SPA tourne.",
        7.75,
        1.55,
        3.8,
        1.95,
    )
    add_code_box(slide, "Avant", ngrx_before, 0.85, 3.85, 5.15, 2.45)
    add_code_box(slide, "Maintenant", ngrx_after, 6.15, 3.85, 5.3, 2.45)
    add_footer(slide, "Vocabulaire cle: withEntities, withState, withComputed, withMethods, withHooks, withStatus")

    forms_before = """readonly profileForm = new FormGroup({
  name: new FormControl('', { nonNullable: true }),
  email: new FormControl('', { nonNullable: true }),
  role: new FormControl('developer', { nonNullable: true })
});"""

    forms_after = """readonly profileModel = signal({
  name: '',
  email: '',
  role: 'developer'
});

readonly profileForm = form(this.profileModel, (path) => {
  required(path.name);
  email(path.email);
});"""

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)
    add_title(slide, "Chapitre 7 - Signal Forms", "Un sujet d ouverture a presenter avec prudence")
    add_section_card(
        slide,
        "Ce que ca change",
        "On part davantage d un modele signal et on construit le formulaire depuis ce modele.",
        0.85,
        1.55,
        3.5,
        1.95,
    )
    add_section_card(
        slide,
        "Impact",
        "Le vocabulaire change: modele, form(), formField et validation rattachee a ce modele.",
        4.5,
        1.55,
        3.15,
        1.95,
    )
    add_section_card(
        slide,
        "Point de vigilance",
        "La documentation Angular le presente encore comme experimental. A montrer comme une direction, pas comme une regle absolue.",
        7.85,
        1.55,
        3.6,
        1.95,
    )
    add_code_box(slide, "Avant", forms_before, 0.85, 3.85, 5.15, 2.1)
    add_code_box(slide, "Maintenant", forms_after, 6.15, 3.85, 5.3, 2.1)
    add_footer(slide, "Message cle: sujet tres formateur, mais encore experimental")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)
    add_title(slide, "Conclusion", "Ce qu on adopte vite, ce qu on cadre, ce qu on surveille")
    add_section_card(
        slide,
        "A adopter vite",
        "Standalone\nNouveau controle de flux\n@let\nSignals\ninput() / output()\nInterop RxJS / Signals",
        0.9,
        1.8,
        3.5,
        3.0,
    )
    add_section_card(
        slide,
        "A cadrer projet",
        "NgRx Signal Store\nParce que le sujet touche a la structure globale de l application et pas seulement a une syntaxe.",
        4.75,
        1.8,
        3.3,
        3.0,
    )
    add_section_card(
        slide,
        "A presenter avec prudence",
        "Signal Forms\nTres interessant pour comprendre la direction du framework, mais encore experimental.",
        8.4,
        1.8,
        3.0,
        3.0,
    )
    add_bullets(
        slide,
        [
            "Message final: Angular recent cherche surtout a rendre le code plus local, plus lisible et plus simple a raisonner.",
        ],
        1.0,
        5.3,
        10.2,
        0.7,
        18,
    )
    add_footer(slide, "Fin de presentation")

    prs.save(output_path)


if __name__ == "__main__":
    project_root = Path(__file__).resolve().parents[1]
    output = project_root / "Decouverte-Angular-21.pptx"
    build_presentation(output)
    print(output)
